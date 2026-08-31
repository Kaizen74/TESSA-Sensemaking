"""File parsers for ingestion (PRD §1.3).

Eight formats in two classes:

* **narrative** — ``.docx`` ``.txt`` ``.md`` ``.pdf`` ``.pptx`` ``.vtt`` ``.srt``.
  Prose, read into a list of text blocks, each carrying the locator a human
  could use to find it again ("page 3, paragraph 2", "slide 4 notes").
* **tabular** — ``.xlsx`` (multi-sheet) ``.csv``. Rows under headers, kept sheet
  by sheet so a mixed-role workbook can be mapped per sheet and "ignore" sheets
  skipped whole (PRD §9 assumption 10).

Parsing is deterministic and offline — no AI is involved at this step, and the
AI never sees the file itself, only the normalised text this module produces.
The file class decides which Stage A prompt runs and which confirmation screen
the operator sees, so it is settled here, once, from the extension.

Decks are read text-only: slide text and speaker notes. Embedded images and
charts are ignored (PRD §9 assumption 9).
"""

from __future__ import annotations

import csv
import io
import re
from pathlib import Path

from pydantic import BaseModel, ConfigDict, Field

FILE_CLASS_NARRATIVE = "narrative"
FILE_CLASS_TABULAR = "tabular"

#: Extension → (file_type, file_class). The single source for both "can this be
#: imported at all" and "which Stage A prompt applies".
FILE_TYPES: dict[str, tuple[str, str]] = {
    ".docx": ("docx", FILE_CLASS_NARRATIVE),
    ".txt": ("txt", FILE_CLASS_NARRATIVE),
    ".md": ("md", FILE_CLASS_NARRATIVE),
    ".pdf": ("pdf", FILE_CLASS_NARRATIVE),
    ".pptx": ("pptx", FILE_CLASS_NARRATIVE),
    ".vtt": ("vtt", FILE_CLASS_NARRATIVE),
    ".srt": ("srt", FILE_CLASS_NARRATIVE),
    ".xlsx": ("xlsx", FILE_CLASS_TABULAR),
    ".csv": ("csv", FILE_CLASS_TABULAR),
}

#: file_type → file_class, for the jobs already in the database, whose file type
#: was recorded but whose extension is long gone.
FILE_CLASSES: dict[str, str] = dict(FILE_TYPES.values())

#: Files larger than this are refused rather than read into memory. A workshop's
#: worth of transcripts is a long way under it.
MAX_UPLOAD_BYTES = 25 * 1024 * 1024

#: Blocks shorter than this are headings, page numbers and stray labels, not
#: stories. They are dropped at parse time so Stage A is not asked about them.
MIN_BLOCK_CHARS = 2


class ParseError(Exception):
    """A file that could not be read, phrased for the operator."""

    def __init__(self, code: str, message: str, action: str) -> None:
        super().__init__(message)
        self.code = code
        self.message = message
        self.action = action


class Block(BaseModel):
    """One piece of prose and where in the file it came from."""

    model_config = ConfigDict(extra="forbid")

    locator: str
    text: str


class Sheet(BaseModel):
    """One sheet of a workbook, or the whole of a CSV."""

    model_config = ConfigDict(extra="forbid")

    name: str
    headers: list[str]
    rows: list[list[str]]
    #: Spreadsheet row number of ``rows[0]``, so a locator points at the row the
    #: operator sees in Excel rather than at a zero-based index.
    first_row_number: int


class NormalisedDocument(BaseModel):
    """A parsed file, in the one shape the rest of ingestion works from."""

    model_config = ConfigDict(extra="forbid")

    file_type: str
    file_class: str
    blocks: list[Block] = Field(default_factory=list)
    sheets: list[Sheet] = Field(default_factory=list)

    @property
    def row_count(self) -> int:
        """Every data row in the file, across every sheet, ignored or not."""
        return sum(len(sheet.rows) for sheet in self.sheets)


def classify(filename: str) -> tuple[str, str]:
    """Return ``(file_type, file_class)`` for a filename, or refuse it."""
    suffix = Path(filename).suffix.lower()
    known = FILE_TYPES.get(suffix)
    if known is None:
        readable = ", ".join(sorted(FILE_TYPES))
        raise ParseError(
            "unsupported_file_type",
            f"Narrative Lens cannot read {suffix or 'files without an extension'}.",
            f"Save the file as one of these and try again: {readable}.",
        )
    return known


def _decode(data: bytes) -> str:
    """Text from bytes, preferring UTF-8 but never failing on an old export."""
    for encoding in ("utf-8-sig", "utf-8", "cp1252"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _blocks_from_text(text: str, prefix: str) -> list[Block]:
    """Split prose on blank lines into located blocks.

    A blank line is the one paragraph marker every plain-text export agrees on,
    and it is what a person writing up thirty workshop responses actually types
    between them.
    """
    blocks: list[Block] = []
    for index, chunk in enumerate(re.split(r"\n\s*\n", text), start=1):
        cleaned = "\n".join(line.rstrip() for line in chunk.strip().splitlines())
        if len(cleaned) >= MIN_BLOCK_CHARS:
            blocks.append(Block(locator=f"{prefix}paragraph {index}", text=cleaned))
    return blocks


def _parse_plain(data: bytes) -> list[Block]:
    return _blocks_from_text(_decode(data), prefix="")


def _parse_docx(data: bytes) -> list[Block]:
    """Paragraphs of a Word document, in document order.

    Text inside Word tables is not read: a ``.docx`` is treated as prose, and a
    table of responses belongs in the tabular path where it gets a confirmed
    column mapping and exact row reconciliation (constraint 12).
    """
    import docx

    try:
        document = docx.Document(io.BytesIO(data))
    except Exception as exc:
        raise ParseError(
            "docx_unreadable",
            "That Word file could not be opened.",
            "Open it in Word and re-save it as .docx, then import it again.",
        ) from exc

    blocks: list[Block] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if len(text) >= MIN_BLOCK_CHARS:
            blocks.append(Block(locator=f"paragraph {index}", text=text))
    return blocks


def _parse_pdf(data: bytes) -> list[Block]:
    """Text of each page, split into paragraphs so locators stay useful."""
    from pypdf import PdfReader

    try:
        reader = PdfReader(io.BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages]
    except Exception as exc:
        raise ParseError(
            "pdf_unreadable",
            "That PDF could not be read.",
            "If it is a scan rather than a document, the words are pictures — "
            "type the stories in under Capture instead.",
        ) from exc

    blocks: list[Block] = []
    for number, text in enumerate(pages, start=1):
        blocks.extend(_blocks_from_text(text, prefix=f"page {number}, "))
    return blocks


def _parse_pptx(data: bytes) -> list[Block]:
    """Slide text and speaker notes, text only (PRD §9 assumption 9)."""
    from pptx import Presentation

    try:
        deck = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise ParseError(
            "pptx_unreadable",
            "That slide deck could not be opened.",
            "Open it in PowerPoint and re-save it as .pptx, then import it again.",
        ) from exc

    blocks: list[Block] = []
    for number, slide in enumerate(deck.slides, start=1):
        parts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        body = "\n".join(parts).strip()
        if len(body) >= MIN_BLOCK_CHARS:
            blocks.append(Block(locator=f"slide {number}", text=body))

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if len(notes) >= MIN_BLOCK_CHARS:
                blocks.append(Block(locator=f"slide {number} notes", text=notes))
    return blocks


#: The timing line both caption formats share. VTT uses dots for fractions,
#: SRT uses commas; nothing else about the line differs.
_STAMP = r"\d{1,2}:\d{2}(?::\d{2})?[.,]\d{1,3}"
CUE_TIMING = re.compile(rf"({_STAMP})\s*-->\s*({_STAMP})")


def _parse_captions(data: bytes) -> list[Block]:
    """Cues of a ``.vtt`` or ``.srt`` transcript, one block each.

    Consecutive cues are kept separate rather than glued into one wall of text:
    Stage A groups them into stories, and a human confirms that grouping. Doing
    it here would be this module deciding where a story begins, which is exactly
    the judgement constraint 1 reserves for a person.
    """
    blocks: list[Block] = []
    number = 0
    for chunk in re.split(r"\n\s*\n", _decode(data)):
        lines = [line.strip() for line in chunk.strip().splitlines() if line.strip()]
        timing = next((line for line in lines if CUE_TIMING.search(line)), None)
        if timing is None:
            continue  # WEBVTT header, NOTE comment, or a stray blank chunk
        match = CUE_TIMING.search(timing)
        assert match is not None  # guarded by the search above
        spoken = "\n".join(lines[lines.index(timing) + 1 :]).strip()
        if len(spoken) < MIN_BLOCK_CHARS:
            continue
        number += 1
        start, end = match.group(1), match.group(2)
        blocks.append(Block(locator=f"cue {number} ({start} → {end})", text=spoken))
    return blocks


def _clean_row(values: list[object]) -> list[str]:
    """Cells as trimmed strings; ``None`` and blanks both become ``''``."""
    return ["" if value is None else str(value).strip() for value in values]


def _sheet_from_rows(name: str, raw_rows: list[list[str]]) -> Sheet | None:
    """Take the first non-empty row as headers and the rest as data."""
    for index, row in enumerate(raw_rows):
        if any(cell for cell in row):
            headers = [cell or f"Column {n}" for n, cell in enumerate(row, start=1)]
            body = [r for r in raw_rows[index + 1 :] if any(cell for cell in r)]
            width = len(headers)
            padded = [(r + [""] * width)[:width] for r in body]
            # +2: one for the header row itself, one because spreadsheet rows
            # are numbered from 1 and this index is from 0.
            return Sheet(
                name=name, headers=headers, rows=padded, first_row_number=index + 2
            )
    return None


def _parse_csv(data: bytes, filename: str) -> list[Sheet]:
    text = _decode(data)
    sample = text[:4096]
    try:
        dialect: type[csv.Dialect] | csv.Dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel  # a single-column file sniffs as nothing; commas are fine
    rows = [_clean_row(list(row)) for row in csv.reader(io.StringIO(text), dialect)]
    sheet = _sheet_from_rows(Path(filename).stem or "Sheet 1", rows)
    return [sheet] if sheet else []


def _parse_xlsx(data: bytes) -> list[Sheet]:
    """Every worksheet of a workbook, kept separate (assumption 10)."""
    from openpyxl import load_workbook

    try:
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        raise ParseError(
            "xlsx_unreadable",
            "That spreadsheet could not be opened.",
            "Open it in Excel and re-save it as .xlsx, then import it again.",
        ) from exc

    try:
        sheets: list[Sheet] = []
        for worksheet in workbook.worksheets:
            rows = [_clean_row(list(row)) for row in worksheet.iter_rows(values_only=True)]
            sheet = _sheet_from_rows(worksheet.title, rows)
            if sheet is not None:
                sheets.append(sheet)
        return sheets
    finally:
        workbook.close()


def parse(filename: str, data: bytes) -> NormalisedDocument:
    """Read one uploaded file into the normalised shape, or refuse it.

    Refusals are plain-English (constraint 7) and always say what to do next,
    because the operator's alternative is never "debug it" — it is "type these
    in under Capture", which still works with no network at all.
    """
    if len(data) > MAX_UPLOAD_BYTES:
        raise ParseError(
            "file_too_large",
            "That file is larger than Narrative Lens will read in one go.",
            f"Split it into parts under {MAX_UPLOAD_BYTES // (1024 * 1024)} MB "
            "and import them one at a time.",
        )

    file_type, file_class = classify(filename)

    blocks: list[Block] = []
    sheets: list[Sheet] = []
    if file_type == "docx":
        blocks = _parse_docx(data)
    elif file_type in ("txt", "md"):
        blocks = _parse_plain(data)
    elif file_type == "pdf":
        blocks = _parse_pdf(data)
    elif file_type == "pptx":
        blocks = _parse_pptx(data)
    elif file_type in ("vtt", "srt"):
        blocks = _parse_captions(data)
    elif file_type == "csv":
        sheets = _parse_csv(data, filename)
    elif file_type == "xlsx":
        sheets = _parse_xlsx(data)

    document = NormalisedDocument(
        file_type=file_type, file_class=file_class, blocks=blocks, sheets=sheets
    )

    if not document.blocks and not document.sheets:
        raise ParseError(
            "file_empty",
            "Narrative Lens opened that file but found no text in it.",
            "Check it is the right file. If it is a scan or a photograph, the "
            "words are pictures — type the stories in under Capture instead.",
        )
    return document
