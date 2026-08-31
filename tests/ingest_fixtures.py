"""Real files, in memory, one per format Narrative Lens claims to read.

Acceptance criterion 7 says all-format fixtures pass Stage A. A fixture that is
a hand-written dict proves nothing about that — the whole risk of ingestion is
in the file formats themselves. So every builder here produces genuine bytes:
a real zip-backed ``.docx``, a real workbook, a real PDF with a real content
stream. The parsers are then held to files of the kind an operator would
actually be handed.

The PDF is written by hand rather than with a rendering library. The app has no
PDF *writer* dependency and PRD §9 assumption 11 is keen it never grows one, so
the test suite builds the smallest file that is genuinely a PDF instead.
"""

from __future__ import annotations

import io

#: Three workshop answers, long enough that the Stage A mock reads them as whole
#: accounts rather than as fragments.
STORY_ONE = (
    "We were three hours from the deadline when the parts finally arrived, and "
    "nobody had told the night shift they were coming. Two of us stayed to "
    "unload them because there was no one else, and the job went out on time."
)
STORY_TWO = (
    "The new checklist looked sensible on paper but it assumed you had both "
    "hands free, which on a wet deck you never do. We ended up doing it from "
    "memory and ticking the boxes afterwards, which is exactly what it was "
    "meant to stop."
)
STORY_THREE = (
    "A customer rang about a fault we had already fixed twice. I could see the "
    "history but not what the engineer had actually done, so I sent someone "
    "out again. She was polite about it. I would not have been."
)

#: Short enough to land under constraint 2's 0.70 in the mock, so the amber path
#: is exercised by a fixture rather than only by a unit test.
SHORT_NOTE = "Ran out of time."


def txt_bytes() -> bytes:
    return "\n\n".join([STORY_ONE, STORY_TWO, SHORT_NOTE]).encode("utf-8")


def md_bytes() -> bytes:
    return ("# Workshop notes\n\n" + "\n\n".join([STORY_ONE, STORY_TWO])).encode("utf-8")


def docx_bytes() -> bytes:
    import docx

    document = docx.Document()
    document.add_paragraph(STORY_ONE)
    document.add_paragraph("")
    document.add_paragraph(STORY_TWO)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "What we heard"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    box.text_frame.text = STORY_ONE
    slide.notes_slide.notes_text_frame.text = STORY_TWO
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def xlsx_bytes() -> bytes:
    """A two-sheet workbook: one of responses, one lookup table to ignore.

    The responses sheet carries one deliberately empty story cell, so the
    reconciliation has something real to account for (constraint 12).
    """
    from openpyxl import Workbook

    workbook = Workbook()
    responses = workbook.active
    responses.title = "Responses"
    responses.append(["Ref", "Team", "Story", "Logged"])
    responses.append(["R-1", "Ops", STORY_ONE, "2026-05-02"])
    responses.append(["R-2", "Deck", STORY_TWO, "2026-05-02"])
    responses.append(["R-3", "Support", "", "2026-05-03"])
    responses.append(["R-4", "Support", STORY_THREE, "2026-05-03"])

    lookup = workbook.create_sheet("Team codes")
    lookup.append(["Code", "Team"])
    lookup.append(["OPS", "Ops"])
    lookup.append(["DCK", "Deck"])

    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def csv_bytes() -> bytes:
    rows = [
        "Team,Story",
        f'Ops,"{STORY_ONE}"',
        f'Deck,"{STORY_TWO}"',
        'Support,""',
    ]
    return ("\n".join(rows) + "\n").encode("utf-8")


def vtt_bytes() -> bytes:
    return (
        "WEBVTT\n\n"
        "00:00:01.000 --> 00:00:12.000\n"
        f"{STORY_ONE}\n\n"
        "00:00:12.500 --> 00:00:24.000\n"
        f"{STORY_TWO}\n"
    ).encode()


def srt_bytes() -> bytes:
    return (
        "1\n"
        "00:00:01,000 --> 00:00:12,000\n"
        f"{STORY_ONE}\n\n"
        "2\n"
        "00:00:12,500 --> 00:00:24,000\n"
        f"{STORY_TWO}\n"
    ).encode()


def _pdf_escape(text: str) -> str:
    return text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")


def pdf_bytes(paragraphs: tuple[str, ...] = (STORY_ONE, STORY_TWO)) -> bytes:
    """A one-page PDF with each paragraph as its own text run.

    Written by hand: catalog, page tree, one page, one Helvetica font, and a
    content stream. Lines are wrapped so a paragraph is several ``Tj`` runs and
    the blank line between paragraphs is a wider leading — which is what makes
    the extracted text split the way a real export would.
    """
    lines: list[str] = []
    for index, paragraph in enumerate(paragraphs):
        if index:
            lines.append("")  # the blank line the parser splits paragraphs on
        words = paragraph.split()
        current = ""
        for word in words:
            if len(current) + len(word) + 1 > 70:
                lines.append(current)
                current = word
            else:
                current = f"{current} {word}".strip()
        if current:
            lines.append(current)

    runs = ["BT", "/F1 11 Tf", "14 TL", "40 780 Td"]
    for line in lines:
        # A blank line is written as a single space: a truly empty Tj run is
        # dropped by every extractor, and then the paragraph break vanishes.
        runs.append(f"({_pdf_escape(line or ' ')}) Tj")
        runs.append("T*")
    runs.append("ET")
    stream = "\n".join(runs).encode("latin-1")

    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{number} 0 obj\n".encode() + body + b"\nendobj\n"

    xref_at = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode()
    out += f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n".encode()
    out += f"{xref_at}\n%%EOF\n".encode()
    return bytes(out)


#: Every format the PRD lists, with a real file for each. Used by the
#: all-formats test that acceptance criterion 7 asks for.
ALL_FORMATS: dict[str, bytes] = {
    "workshop.txt": txt_bytes(),
    "workshop.md": md_bytes(),
    "workshop.docx": docx_bytes(),
    "workshop.pdf": pdf_bytes(),
    "workshop.pptx": pptx_bytes(),
    "workshop.xlsx": xlsx_bytes(),
    "workshop.csv": csv_bytes(),
    "workshop.vtt": vtt_bytes(),
    "workshop.srt": srt_bytes(),
}
